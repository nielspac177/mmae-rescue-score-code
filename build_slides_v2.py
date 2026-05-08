"""SVIN / grand-rounds slide deck for the MMA-embolization rescue score.
JAMA Neurology aesthetic, 16:9, 10 slides."""
from __future__ import annotations
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

HERE = Path(__file__).parent
V2 = HERE / "v2"
OUT = HERE / "MMA_rescue_slides.pptx"

# JAMA Neurology palette
NAVY    = RGBColor(0x37, 0x4E, 0x55)
NAVY_DK = RGBColor(0x2A, 0x3C, 0x40)
GOLD    = RGBColor(0xDF, 0x8F, 0x44)
GOLD_DK = RGBColor(0xB9, 0x85, 0x38)
PURPLE  = RGBColor(0x6A, 0x65, 0x99)
RED     = RGBColor(0xB2, 0x47, 0x45)
GREEN   = RGBColor(0x79, 0xAF, 0x97)
GREY    = RGBColor(0x80, 0x79, 0x6B)
PALE    = RGBColor(0xF5, 0xF2, 0xEA)
INK     = RGBColor(0x22, 0x22, 0x22)
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
SAND    = RGBColor(0xEF, 0xE9, 0xDB)


def add_text(slide, text, left, top, width, height,
             font_size=18, bold=False, color=INK, align="left",
             font="Helvetica"):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0); tf.margin_right = Emu(0)
    tf.margin_top  = Emu(0); tf.margin_bottom = Emu(0)
    p = tf.paragraphs[0]
    p.alignment = {"left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER,
                   "right": PP_ALIGN.RIGHT}[align]
    r = p.add_run()
    r.text = text
    r.font.name = font
    r.font.size = Pt(font_size)
    r.font.bold = bold
    r.font.color.rgb = color
    return tb


def add_bullets(slide, items, left, top, width, height,
                font_size=18, color=INK, font="Helvetica"):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    for i, txt in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(6)
        r = p.add_run()
        r.text = "• " + txt
        r.font.name = font
        r.font.size = Pt(font_size)
        r.font.color.rgb = color
    return tb


def add_rect(slide, left, top, width, height, fill=NAVY, line=None):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
        shp.line.width = Pt(0.5)
    shp.shadow.inherit = False
    return shp


def add_footer(slide, slide_no, total):
    add_rect(slide, Inches(0), Inches(7.10), Inches(13.33), Inches(0.40),
             fill=NAVY)
    add_text(slide, "MMA Embolization Rescue Score  ·  Pacheco et al., 2026",
             Inches(0.4), Inches(7.18), Inches(8.0), Inches(0.30),
             font_size=10.5, color=WHITE)
    add_text(slide, f"{slide_no} / {total}",
             Inches(11.8), Inches(7.18), Inches(1.2), Inches(0.30),
             font_size=10.5, color=WHITE, align="right")


# ----------------------------------------------------------------------
def main():
    prs = Presentation()
    prs.slide_width  = Inches(13.33)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]   # blank
    TOTAL = 10

    # ============================================================
    # SLIDE 1 — TITLE
    # ============================================================
    s = prs.slides.add_slide(blank)
    # navy banner
    add_rect(s, Inches(0), Inches(0), Inches(13.33), Inches(7.5), fill=NAVY)
    add_rect(s, Inches(0), Inches(2.25), Inches(13.33), Inches(0.05),
             fill=GOLD)
    add_rect(s, Inches(0), Inches(4.85), Inches(13.33), Inches(0.05),
             fill=GOLD)

    add_text(s, "A 7-point bedside score for rescue surgery",
             Inches(0.7), Inches(2.55), Inches(12), Inches(0.85),
             font_size=40, bold=True, color=WHITE)
    add_text(s, "after middle meningeal artery embolization for cSDH",
             Inches(0.7), Inches(3.40), Inches(12), Inches(0.7),
             font_size=28, color=GOLD)

    add_text(s, "Single-center derivation and internal validation, n = 214",
             Inches(0.7), Inches(4.15), Inches(12), Inches(0.5),
             font_size=18, color=WHITE)

    add_text(s, "Pacheco-Barrios N, Gonzales-Salidos J, et al.",
             Inches(0.7), Inches(5.10), Inches(12), Inches(0.4),
             font_size=18, color=WHITE)
    add_text(s, "Department of Neurosurgery  ·  2026",
             Inches(0.7), Inches(5.55), Inches(12), Inches(0.4),
             font_size=14, color=GOLD)

    add_text(s, "nielspac177.github.io/mmae-rescue-score",
             Inches(0.7), Inches(6.60), Inches(12), Inches(0.4),
             font_size=14, color=GOLD, font="DejaVu Sans Mono")

    # ============================================================
    # SLIDE 2 — BACKGROUND
    # ============================================================
    s = prs.slides.add_slide(blank)
    add_rect(s, Inches(0), Inches(0), Inches(13.33), Inches(0.7), fill=NAVY)
    add_text(s, "Background", Inches(0.5), Inches(0.15), Inches(12), Inches(0.5),
             font_size=24, bold=True, color=WHITE)

    add_text(s, "Three randomized trials in 2024–2025 cemented MMA embolization as standard adjunct for cSDH.",
             Inches(0.5), Inches(1.0), Inches(12.3), Inches(0.6),
             font_size=20, bold=True, color=NAVY)

    add_bullets(s, [
        "EMBOLISE (N=400): 4.1% vs 11.3% reoperation at 90 days (Kan 2024)",
        "MAGIC-MT (N=722): same direction of effect (Tian 2024)",
        "EMPROTECT (N=342): significant 6-month recurrence reduction (Kerleroux 2024)",
    ], Inches(0.7), Inches(1.7), Inches(12.0), Inches(2.0),
       font_size=17, color=INK)

    add_text(s, "But 5–20% of embolized patients still need rescue surgery.",
             Inches(0.5), Inches(3.85), Inches(12.3), Inches(0.6),
             font_size=20, bold=True, color=RED)

    add_bullets(s, [
        "No bedside-usable, externally validated score exists",
        "Existing cSDH grades (Markwalder, Stanišić) predict recurrence after burr-hole, not after embolization",
        "Decision-relevant question: who needs tighter post-procedure surveillance?",
    ], Inches(0.7), Inches(4.55), Inches(12.0), Inches(2.0),
       font_size=17, color=INK)

    add_footer(s, 2, TOTAL)

    # ============================================================
    # SLIDE 3 — METHODS
    # ============================================================
    s = prs.slides.add_slide(blank)
    add_rect(s, Inches(0), Inches(0), Inches(13.33), Inches(0.7), fill=NAVY)
    add_text(s, "Methods — two-tier prediction-model strategy",
             Inches(0.5), Inches(0.15), Inches(12), Inches(0.5),
             font_size=24, bold=True, color=WHITE)

    # left column: cohort
    add_rect(s, Inches(0.4), Inches(1.0), Inches(6.2), Inches(0.05),
             fill=NAVY)
    add_text(s, "Cohort", Inches(0.4), Inches(1.1), Inches(6.2), Inches(0.4),
             font_size=18, bold=True, color=NAVY)
    add_bullets(s, [
        "214 consecutive adults, MMA embolization for cSDH",
        "36 rescue events (16.8%); single center, retrospective",
        "TRIPOD-AI compliant; IRB approved with consent waiver",
    ], Inches(0.45), Inches(1.55), Inches(6.2), Inches(2.0),
       font_size=15, color=INK)

    # left column: primary models
    add_rect(s, Inches(0.4), Inches(3.50), Inches(6.2), Inches(0.05),
             fill=NAVY)
    add_text(s, "Primary models — knowledge-driven",
             Inches(0.4), Inches(3.60), Inches(6.2), Inches(0.4),
             font_size=18, bold=True, color=NAVY)
    add_bullets(s, [
        "Model 1 (max 7): age, SDH ≥100 mL, anticoag, plt <150, antiplatelet, ant+post",
        "Model 2 (max 4): age, SDH ≥100 mL, anticoag",
        "Capped at EPV-5 (≤7 predictors per Peduzzi 1996)",
    ], Inches(0.45), Inches(4.05), Inches(6.2), Inches(2.5),
       font_size=15, color=INK)

    # right column: sensitivity models + analysis
    add_rect(s, Inches(6.95), Inches(1.0), Inches(6.0), Inches(0.05),
             fill=GOLD_DK)
    add_text(s, "Sensitivity models",
             Inches(6.95), Inches(1.1), Inches(6.0), Inches(0.4),
             font_size=18, bold=True, color=GOLD_DK)
    add_bullets(s, [
        "Model 3 — alternative age-stratification (age >85, plt, antiplatelet)",
        "Model 4 — lasso, EPV-5 cap, 3 imbalance modes",
        "Model 5 — splines + interactions + tuned EN + stacking + Bayesian + MICE",
        "Focal-deficit-included variants of M1–M3",
    ], Inches(7.0), Inches(1.55), Inches(6.0), Inches(2.5),
       font_size=15, color=INK)

    add_rect(s, Inches(6.95), Inches(3.50), Inches(6.0), Inches(0.05),
             fill=GOLD_DK)
    add_text(s, "Validation",
             Inches(6.95), Inches(3.60), Inches(6.0), Inches(0.4),
             font_size=18, bold=True, color=GOLD_DK)
    add_bullets(s, [
        "Apparent + Harrell optimism-corrected AUC (1000 bootstrap)",
        "Hosmer–Lemeshow calibration",
        "Decision curve analysis (Vickers net-benefit)",
        "5×10 stratified CV for ML/data-driven models",
    ], Inches(7.0), Inches(4.05), Inches(6.0), Inches(2.5),
       font_size=15, color=INK)

    add_footer(s, 3, TOTAL)

    # ============================================================
    # SLIDE 4 — COHORT
    # ============================================================
    s = prs.slides.add_slide(blank)
    add_rect(s, Inches(0), Inches(0), Inches(13.33), Inches(0.7), fill=NAVY)
    add_text(s, "Cohort characteristics",
             Inches(0.5), Inches(0.15), Inches(12), Inches(0.5),
             font_size=24, bold=True, color=WHITE)

    rows = [
        ("Variable",                       "Cohort (n=214)"),
        ("Age, mean (SD), years",          "73.4 (13.1)"),
        ("Age >80 years",                  "31.8%"),
        ("Anticoagulation",                "25.7%"),
        ("Antiplatelet therapy",           "36.4%"),
        ("SDH volume ≥100 mL",             "29.0%"),
        ("Platelets <150 ×10⁹/L",          "20.6%"),
        ("Anterior + posterior embolization", "53.7%"),
        ("Focal deficit at presentation",  "59.8%"),
        ("Rescue surgery (primary outcome)", "36 (16.8%)"),
    ]
    table_left = Inches(2.0); table_top = Inches(1.1)
    col_w = [Inches(6.5), Inches(3.5)]
    row_h = Inches(0.45)
    for i, (a, b) in enumerate(rows):
        is_header = i == 0
        is_outcome = i == len(rows) - 1
        bg = NAVY if is_header else (SAND if is_outcome else
                                       (PALE if i % 2 == 0 else WHITE))
        for j, (txt, w) in enumerate(zip([a, b], col_w)):
            x = table_left + (Inches(0) if j == 0 else col_w[0])
            y = table_top + row_h * i
            add_rect(s, x, y, w, row_h, fill=bg, line=GREY)
            txt_color = WHITE if is_header else (NAVY if is_outcome else INK)
            add_text(s, txt, x + Inches(0.15), y + Inches(0.06),
                     w - Inches(0.3), row_h - Inches(0.1),
                     font_size=14,
                     bold=is_header or is_outcome,
                     color=txt_color)

    add_footer(s, 4, TOTAL)

    # ============================================================
    # SLIDE 5 — PRIMARY RESULT (ROC + Score-rate)
    # ============================================================
    s = prs.slides.add_slide(blank)
    add_rect(s, Inches(0), Inches(0), Inches(13.33), Inches(0.7), fill=NAVY)
    add_text(s, "Primary result — Model 1 (knowledge-driven, 7 pts)",
             Inches(0.5), Inches(0.15), Inches(12), Inches(0.5),
             font_size=24, bold=True, color=WHITE)

    # Headline AUC
    add_text(s, "AUC 0.706  (Harrell optimism-corrected 0.704)",
             Inches(0.5), Inches(0.95), Inches(12), Inches(0.45),
             font_size=22, bold=True, color=NAVY)
    add_text(s, "Hosmer–Lemeshow P > 0.5  ·  Brier 0.13  ·  TRIPOD-AI compliant",
             Inches(0.5), Inches(1.45), Inches(12), Inches(0.4),
             font_size=14, color=GREY)

    # ROC on left, score-rate on right
    s.shapes.add_picture(str(V2 / "fig1_roc.png"),
                         Inches(0.5), Inches(2.0),
                         height=Inches(4.6))
    s.shapes.add_picture(str(V2 / "fig2_score_risk.png"),
                         Inches(5.7), Inches(2.6),
                         width=Inches(7.5))

    add_text(s, "Stepwise rise in observed rescue rate across the 0–6 score range",
             Inches(5.7), Inches(6.55), Inches(7.5), Inches(0.4),
             font_size=12, color=GREY, align="center")

    add_footer(s, 5, TOTAL)

    # ============================================================
    # SLIDE 6 — DICHOTOMIZATION
    # ============================================================
    s = prs.slides.add_slide(blank)
    add_rect(s, Inches(0), Inches(0), Inches(13.33), Inches(0.7), fill=NAVY)
    add_text(s, "A single bedside cutoff: score ≥ 4",
             Inches(0.5), Inches(0.15), Inches(12), Inches(0.5),
             font_size=24, bold=True, color=WHITE)

    # Big numbers: 8.6% vs 36.5% with arrow / 4.2× chip
    add_text(s, "Low risk", Inches(1.0), Inches(1.2), Inches(4.5), Inches(0.5),
             font_size=22, bold=True, color=GREEN, align="center")
    add_text(s, "Score 0–3", Inches(1.0), Inches(1.65), Inches(4.5), Inches(0.4),
             font_size=14, color=GREY, align="center")
    add_text(s, "8.6%", Inches(1.0), Inches(2.15), Inches(4.5), Inches(1.1),
             font_size=72, bold=True, color=GREEN, align="center")
    add_text(s, "13 of 151 patients",
             Inches(1.0), Inches(3.30), Inches(4.5), Inches(0.4),
             font_size=14, color=GREY, align="center")

    add_text(s, "High risk", Inches(7.85), Inches(1.2), Inches(4.5), Inches(0.5),
             font_size=22, bold=True, color=RED, align="center")
    add_text(s, "Score ≥ 4", Inches(7.85), Inches(1.65), Inches(4.5), Inches(0.4),
             font_size=14, color=GREY, align="center")
    add_text(s, "36.5%", Inches(7.85), Inches(2.15), Inches(4.5), Inches(1.1),
             font_size=72, bold=True, color=RED, align="center")
    add_text(s, "23 of 63 patients",
             Inches(7.85), Inches(3.30), Inches(4.5), Inches(0.4),
             font_size=14, color=GREY, align="center")

    # Center 4.2x ratio chip
    add_rect(s, Inches(5.9), Inches(2.4), Inches(1.55), Inches(0.95),
             fill=NAVY)
    add_text(s, "4.2×", Inches(5.9), Inches(2.45), Inches(1.55), Inches(0.7),
             font_size=32, bold=True, color=GOLD, align="center")
    add_text(s, "rescue gradient",
             Inches(5.9), Inches(3.40), Inches(1.55), Inches(0.3),
             font_size=10, color=NAVY, align="center")

    # Operating-point summary
    add_text(s, "Operating point",
             Inches(0.5), Inches(4.20), Inches(12), Inches(0.4),
             font_size=18, bold=True, color=NAVY)
    add_bullets(s, [
        "Sensitivity 64% (95% CI 48–78)  ·  Specificity 78% (71–83)",
        "Positive predictive value 37%  ·  Negative predictive value 91%",
        "Decision-curve analysis: positive net benefit across 10–40% threshold range",
        "Bedside use: PDF score card and browser calculator (no PHI leaves device)",
    ], Inches(0.7), Inches(4.65), Inches(12.0), Inches(2.5),
       font_size=15, color=INK)

    add_footer(s, 6, TOTAL)

    # ============================================================
    # SLIDE 7 — CALIBRATION + DCA
    # ============================================================
    s = prs.slides.add_slide(blank)
    add_rect(s, Inches(0), Inches(0), Inches(13.33), Inches(0.7), fill=NAVY)
    add_text(s, "Calibration and decision-curve analysis",
             Inches(0.5), Inches(0.15), Inches(12), Inches(0.5),
             font_size=24, bold=True, color=WHITE)

    s.shapes.add_picture(str(V2 / "fig3_calibration.png"),
                         Inches(0.5), Inches(1.0),
                         height=Inches(5.6))
    s.shapes.add_picture(str(V2 / "fig10_dca.png"),
                         Inches(6.7), Inches(1.2),
                         height=Inches(5.4))

    add_text(s, "Predicted-vs-observed across deciles  ·  HL P > 0.5",
             Inches(0.5), Inches(6.65), Inches(5.7), Inches(0.4),
             font_size=12, color=GREY, align="center")
    add_text(s, "Net benefit > treat-all and treat-none across 10–40% threshold range",
             Inches(6.7), Inches(6.65), Inches(6.5), Inches(0.4),
             font_size=12, color=GREY, align="center")

    add_footer(s, 7, TOTAL)

    # ============================================================
    # SLIDE 8 — SENSITIVITY
    # ============================================================
    s = prs.slides.add_slide(blank)
    add_rect(s, Inches(0), Inches(0), Inches(13.33), Inches(0.7), fill=NAVY)
    add_text(s, "Sensitivity analyses — discrimination preserved",
             Inches(0.5), Inches(0.15), Inches(12), Inches(0.5),
             font_size=24, bold=True, color=WHITE)

    add_text(s, "AUCs across all five model families overlap the Model 1 95% CI",
             Inches(0.5), Inches(1.0), Inches(12.3), Inches(0.5),
             font_size=18, color=NAVY)

    # ML benchmark figure
    s.shapes.add_picture(str(V2 / "fig8_ml_bars.png"),
                         Inches(0.5), Inches(1.6),
                         width=Inches(7.5))

    # Right side — sensitivity summary
    add_rect(s, Inches(8.30), Inches(1.6), Inches(4.65), Inches(0.05),
             fill=GOLD_DK)
    add_text(s, "Sensitivity AUC summary",
             Inches(8.30), Inches(1.70), Inches(4.65), Inches(0.4),
             font_size=16, bold=True, color=GOLD_DK)
    add_bullets(s, [
        "Model 2 (simple, 4 pts)  AUC 0.64",
        "Model 3 (age cutoff >85, 5 pts)  AUC 0.70",
        "Model 4 (lasso, 5 pts)  CV-AUC 0.66",
        "Model 5 (Bayesian + priors)  CV-AUC 0.67",
        "Composite enhanced bundle  CV-AUC 0.63",
        "Random forest / GBM / XGB  CV-AUC 0.57–0.58",
    ], Inches(8.30), Inches(2.20), Inches(4.65), Inches(4.2),
       font_size=13, color=INK)

    add_text(s, "36 events bound the AUC ceiling — flexibility costs more than it gains.",
             Inches(0.5), Inches(6.35), Inches(12.3), Inches(0.5),
             font_size=15, color=NAVY, bold=True)

    add_footer(s, 8, TOTAL)

    # ============================================================
    # SLIDE 9 — LIMITATIONS
    # ============================================================
    s = prs.slides.add_slide(blank)
    add_rect(s, Inches(0), Inches(0), Inches(13.33), Inches(0.7), fill=NAVY)
    add_text(s, "Limitations and roadmap",
             Inches(0.5), Inches(0.15), Inches(12), Inches(0.5),
             font_size=24, bold=True, color=WHITE)

    add_text(s, "Limitations",
             Inches(0.5), Inches(1.0), Inches(6.0), Inches(0.4),
             font_size=20, bold=True, color=NAVY)
    add_bullets(s, [
        "Single-center retrospective design",
        "36 events — wide CIs on individual coefficients",
        "Rescue ascertained at last follow-up, not fixed time point",
        "15% of SDH volumes imputed (sensitivity supports primary)",
        "Embolic agent and access route captured but not in score",
    ], Inches(0.6), Inches(1.45), Inches(6.0), Inches(4.0),
       font_size=15, color=INK)

    add_text(s, "Validation roadmap",
             Inches(7.0), Inches(1.0), Inches(6.0), Inches(0.4),
             font_size=20, bold=True, color=GOLD_DK)
    add_bullets(s, [
        "Pre-registration of external-validation analysis plan",
        "Pooled validation across active MMA registries",
        "Embed scoring within next-generation RCTs (EMBOLISE+, EMPROTECT-2)",
        "Update calibration with hierarchical Bayesian recalibration",
        "Add prospective surveillance time-to-rescue endpoint",
    ], Inches(7.1), Inches(1.45), Inches(6.0), Inches(4.0),
       font_size=15, color=INK)

    add_footer(s, 9, TOTAL)

    # ============================================================
    # SLIDE 10 — TAKE HOME
    # ============================================================
    s = prs.slides.add_slide(blank)
    add_rect(s, Inches(0), Inches(0), Inches(13.33), Inches(7.5), fill=NAVY)
    add_rect(s, Inches(0), Inches(2.10), Inches(13.33), Inches(0.05),
             fill=GOLD)

    add_text(s, "Take-home", Inches(0.7), Inches(0.6),
             Inches(12), Inches(0.6),
             font_size=28, bold=True, color=GOLD)

    add_text(s, "A 7-point pre-procedural score, derived from six routinely available variables, "
                "discriminates rescue surgery after MMA embolization for cSDH with AUC 0.70 "
                "and stratifies the cohort into a 4.2-fold rescue-rate gradient at a single cutoff.",
             Inches(0.7), Inches(1.10), Inches(12), Inches(1.0),
             font_size=18, color=WHITE)

    # Three-column tile
    tiles = [
        ("Clinical use",
         "Pre-procedural triage of post-embolization surveillance intensity",
         GREEN),
        ("Methodology",
         "Knowledge-driven primary, lasso + Bayesian + ML as sensitivity",
         GOLD),
        ("Open science",
         "Code, calculator, bedside card, and supplementary tables freely available",
         PURPLE),
    ]
    for i, (title, body, color) in enumerate(tiles):
        x = Inches(0.7 + i * 4.2)
        add_rect(s, x, Inches(2.6), Inches(3.95), Inches(2.6),
                 fill=NAVY_DK)
        add_rect(s, x, Inches(2.6), Inches(3.95), Inches(0.06),
                 fill=color)
        add_text(s, title, x + Inches(0.25), Inches(2.85),
                 Inches(3.45), Inches(0.5),
                 font_size=18, bold=True, color=color)
        add_text(s, body, x + Inches(0.25), Inches(3.40),
                 Inches(3.45), Inches(2.0),
                 font_size=14, color=WHITE)

    # Calculator URL
    add_rect(s, Inches(0.7), Inches(5.45), Inches(12), Inches(1.0),
             fill=NAVY_DK)
    add_text(s, "Live calculator + printable bedside card",
             Inches(0.95), Inches(5.55), Inches(12), Inches(0.4),
             font_size=14, color=GOLD)
    add_text(s, "nielspac177.github.io/mmae-rescue-score",
             Inches(0.95), Inches(5.95), Inches(12), Inches(0.5),
             font_size=22, bold=True, color=WHITE,
             font="DejaVu Sans Mono")

    add_text(s, "Pacheco-Barrios N, Gonzales-Salidos J, et al.  ·  TRIPOD-AI compliant  ·  Pending external validation",
             Inches(0.7), Inches(6.95), Inches(12), Inches(0.4),
             font_size=11, color=GOLD)

    # Save
    prs.save(str(OUT))
    print(f"Wrote {OUT.relative_to(HERE)}  "
          f"({OUT.stat().st_size/1024:.0f} KB, {TOTAL} slides)")


if __name__ == "__main__":
    main()
